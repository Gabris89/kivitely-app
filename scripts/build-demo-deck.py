# -*- coding: utf-8 -*-
"""
A Kivitely chapter-demo prezentacio ujraepitese.

Miert szkript es nem kezi szerkesztes: a kepek a capture-screenshots.mjs
futtatasaval barmikor frissulnek a FUTO appbol, es a dia-sorrend is valtozik
meg. Igy a prezentacio ujragenerelheto, nem kell kezzel utananyulni.

FONTOS korlat: a dia-SZOVEGEK itt vannak beegetve. A szkript nem olvassa a
kodbazist, tehat egy uj funkcio nem jelenik meg magatol - ide kell beirni.
Ami automatikusan kovveti a kodot: a kepernyokepek.

Terv:
  - kevesebb, levegosebb dia (a reszletek a jegyzetbe kerulnek)
  - a kepek nagyok es kerettel elvalnak a hatterrol
  - a hatter SOTETEBB, mint az app sajat hattere, kulonben a kepernyokep
    beleolvad es nem latszik, hol kezdodik

Hasznalat:  python scripts/build-demo-deck.py
Eredmeny:   docs/kivitely-chapter-demo.pptx
"""

import os
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
SHOTS = ROOT / "docs" / "screenshots"
# A kimenet felulirhato: DECK_OUTPUT=... python scripts/build-demo-deck.py
# (hasznos, ha az eredeti fajl epp meg van nyitva PowerPointban)
OUTPUT = Path(os.environ.get("DECK_OUTPUT") or (ROOT / "docs" / "kivitely-chapter-demo.pptx"))

SLIDE_W, SLIDE_H = Inches(17.78), Inches(10.0)

ACCENT = RGBColor(0x68, 0xE1, 0xFD)   # marka-cian (ugyanaz, mint az appban)
BG = RGBColor(0x05, 0x0A, 0x14)       # SOTETEBB, mint az app #0E1A2F hattere
FRAME = RGBColor(0x3A, 0x4E, 0x6B)    # keret a kepernyokepek korul
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY = RGBColor(0xD8, 0xE7, 0xFF)
MUTED = RGBColor(0xAF, 0xC4, 0xE4)
FAINT = RGBColor(0x8F, 0xA6, 0xC5)
FONT = "Aptos"


def textbox(slide, left, top, width, height, text, size, color,
            bold=False, align=PP_ALIGN.LEFT, spacing=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True

    for i, line in enumerate(text if isinstance(text, list) else [text]):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.alignment = align
        if spacing:
            para.space_after = Pt(spacing)
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT
    return box


def note(slide, text):
    """A narracio, amit CSAK az eloado lat (PowerPoint jegyzet resz)."""
    slide.notes_slide.notes_text_frame.text = text


def chrome(prs, eyebrow, title, compact=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bar = slide.shapes.add_shape(1, 0, 0, Inches(0.39), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    bg = slide.shapes.add_shape(1, Inches(0.39), 0, Inches(17.39), SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()

    textbox(slide, 1.0, 0.45, 15.0, 0.24, eyebrow, 14, ACCENT, bold=True)
    if title:
        # Kepes dian kisebb a cim, hogy tobb hely maradjon a kepnek.
        textbox(slide, 1.0, 0.95, 15.5, 0.6, title, 30 if compact else 38, WHITE, bold=True)
    return slide


def bullets(prs, eyebrow, title, items, lead=None, narration=None):
    slide = chrome(prs, eyebrow, title)
    top = 2.3

    if lead:
        textbox(slide, 1.0, top, 15.5, 0.5, lead, 22, MUTED)
        top += 0.9

    textbox(slide, 1.33, top, 15.0, 6.0, [f"• {b}" for b in items], 25, BODY, spacing=18)
    if narration:
        note(slide, narration)
    return slide


def _fit(path, max_w, max_h):
    with Image.open(path) as img:
        w, h = img.size
    scale = min(max_w / w, max_h / h)
    return w * scale, h * scale


def _framed(slide, path, left, top, width, height):
    """Kep vekony kerettel - enelkul a sotet app-kepernyokep beleolvad a hatterbe."""
    pic = slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(width), Inches(height))
    pic.line.color.rgb = FRAME
    pic.line.width = Pt(1.25)
    return pic


def picture(prs, eyebrow, title, filename, caption=None, narration=None):
    slide = chrome(prs, eyebrow, title, compact=True)
    path = SHOTS / filename

    if not path.exists():
        textbox(slide, 1.33, 3.0, 14.0, 1.0, f"[hianyzo kep: {filename}]", 20, FAINT)
        return slide

    top = 1.8
    max_h = 6.9 if caption else 7.4
    w, h = _fit(path, 16.2, max_h)
    _framed(slide, path, (17.78 - w) / 2, top, w, h)

    if caption:
        textbox(slide, 1.0, top + h + 0.2, 15.78, 0.4, caption, 18, MUTED, align=PP_ALIGN.CENTER)
    if narration:
        note(slide, narration)
    return slide


def compare(prs, eyebrow, title, items, caption=None, narration=None):
    slide = chrome(prs, eyebrow, title, compact=True)
    count = len(items)
    gap = 0.35
    total_w = 16.2
    cell_w = (total_w - gap * (count - 1)) / count
    label_top, img_top = 1.85, 2.3
    max_h = 5.6

    for i, (filename, label) in enumerate(items):
        left = (17.78 - total_w) / 2 + i * (cell_w + gap)
        textbox(slide, left, label_top, cell_w, 0.35, label, 17, ACCENT, bold=True, align=PP_ALIGN.CENTER)

        path = SHOTS / filename
        if not path.exists():
            textbox(slide, left, img_top, cell_w, 0.5, f"[hianyzo: {filename}]", 14, FAINT)
            continue

        w, h = _fit(path, cell_w, max_h)
        _framed(slide, path, left + (cell_w - w) / 2, img_top, w, h)

    if caption:
        textbox(slide, 1.0, img_top + max_h + 0.35, 15.78, 0.6, caption, 21, BODY, align=PP_ALIGN.CENTER)
    if narration:
        note(slide, narration)
    return slide


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    # ── 1. Cim ──────────────────────────────────────────────────────────────
    s = chrome(prs, "AI-ASSISTED MVP DEVELOPMENT", None)
    textbox(s, 1.0, 1.19, 14.58, 0.91, "Kivitely", 54, WHITE, bold=True)
    textbox(s, 1.06, 2.2, 14.0, 0.4, "Hogyan indítottunk el egy építőipari MVP-t AI segítségével", 26, MUTED)
    textbox(s, 1.33, 3.3, 14.44, 2.0,
            ["• Domain-specifikus workflow app", "• Next.js + Supabase", "• Valós mobil tesztelés"],
            25, BODY, spacing=18)
    note(s, "Rövid bemutatkozás. A lényeg amit elöljáróban érdemes leszögezni: ez nem egy 'egy prompt "
             "és kész' demó. Egy valódi, több hetes iteráció, ahol a termékdöntések végig emberi kézben "
             "maradtak, az AI pedig implementációs partner volt.")

    # ── 2. Problema ─────────────────────────────────────────────────────────
    bullets(prs, "MIÉRT INDULT?", "Kiinduló probléma",
            ["A terepi információ chatben, fotókban és Excelben szóródik szét.",
             "Hiba, akadály, bizonyítás és elszámolás nincs egy folyamatban.",
             "Cél: terepen is használható munkafolyamat, nem irodai adminfelület."],
            narration="Itt érdemes egy konkrét példát mondani: egy hibát ma lefotóznak telefonnal, "
                      "elküldik WhatsAppon, valaki beírja Excelbe, aztán a hónap végi elszámolásnál "
                      "senki nem találja meg, hogy az adott tétel el lett-e fogadva. Ez a fájdalom.")

    # ── 3. Szerepek ─────────────────────────────────────────────────────────
    bullets(prs, "KINEK SZÓL?", "Öt szerep, öt különböző nap",
            ["Projektvezető – több projekt, pénzügyi rálátás, TIG jóváhagyás.",
             "Építésvezető – egy projekt napi irányítása, kiosztás, átvétel.",
             "Alvállalkozó – a saját munkája: javít, fotóval bizonyít, akadályt jelent.",
             "Megrendelő – betekintés írás nélkül. Adminisztrátor – törzsadatok."],
            lead="A jogosultsági rendszer később pontosan ebből a listából következik.",
            narration="Ez a dia készíti elő a jogosultsági blokkot. Hangsúlyozd, hogy ezek nem kitalált "
                      "szerepek: mindegyik mögött valódi ember áll, akinek MÁS a napja. Az alvállalkozó "
                      "nem akar más cégek hibáit látni, a megrendelő pedig nem akar véletlenül elrontani "
                      "semmit. A jogosultság nem technikai dísz, hanem ebből a valóságból jön.")

    # ── 4. Mi a Kivitely ────────────────────────────────────────────────────
    bullets(prs, "AKTUÁLIS MVP", "Mi a Kivitely ma?",
            ["Projektalapú navigáció, hibalista, workflow tábla, fotós bizonyítás.",
             "Akadálylista, teljesítménynapló, alvállalkozói törzsadat.",
             "Tervek feltöltése és kalibrálható tervmérés.",
             "TIG-előkészítés, Excel és PDF export fotós melléklettel."],
            narration="Gyors felsorolás, ne időzz rajta – a következő diák úgyis megmutatják. "
                      "Amit érdemes kiemelni: ez már nem kattintható prototípus, hanem valódi "
                      "adatbázissal, bejelentkezéssel és jogosultsággal működő alkalmazás.")

    # ── 5-7. Kepek: a termek ────────────────────────────────────────────────
    picture(prs, "BELÉPÉSI PONT", "Áttekintés: minden projekt egy képernyőn",
            "01-attekintes-dashboard.png",
            caption="A nap első kérdése: hol égnek a dolgok?",
            narration="Ez a kezdőképernyő. Aggregált számok minden projektről, alatta állapot szerinti "
                      "eloszlás. Fontos: amit itt látsz, az a BEJELENTKEZETT szerep hatóköre – egy "
                      "alvállalkozó ugyanezen az oldalon jóval kevesebb számot lát.")

    picture(prs, "A TERMÉK MAGJA", "Hibalista: ez váltja ki az Excelt",
            "04-hibalista.png",
            caption="Kereshető, státusz szerint szűrhető, projektenként bontható.",
            narration="A hibalista a napi munka gerince. Minden sor egy elvégzendő feladat: mi a hiba, "
                      "melyik projektben, melyik alvállalkozó csinálja, mikorra. A státusz-gombok fölötte "
                      "egyben szűrők és számlálók.")

    picture(prs, "HIBA RÉSZLETEI", "Egy hiba teljes története",
            "05-hiba-reszletei-allapot-utvonal.png",
            caption="Állapot-útvonal, felelős, határidő, fotós bizonyíték, TIG-készültség.",
            narration="Itt fut össze minden. Felül az állapot-útvonal mutatja, hol tart a hiba az "
                      "életciklusában. Alul a fotók: előtte-utána bizonyítás. A jobb oldali jelzés azt "
                      "mondja meg, mi hiányzik még ahhoz, hogy ez a tétel elszámolható legyen.")

    # ── 8. Osszefuggesek ────────────────────────────────────────────────────
    bullets(prs, "ÖSSZEFÜGGÉSEK", "Hogyan kapcsolódnak az adatok?",
            ["Projekt a gyökér – minden más hozzá tartozik.",
             "Hiba → alvállalkozó (ki csinálja) + fotók (elkészült-e).",
             "Akadály: ami MIATT nem lehet dolgozni – külön él a hibától.",
             "TIG csomag: elfogadott hibák + fotóik = a számla alapja."],
            lead="Egyetlen szabály tartja össze: minden a projekthez tartozik, és minden nyomon követhető.",
            narration="Ez a dia a domain modell lelke. A legfontosabb megkülönböztetés, amit sokan "
                      "összekevernek: a HIBA az, amit el kell végezni, az AKADÁLY viszont az, ami "
                      "megakadályozza a munkát – például nem érkezett meg az anyag. A kettő külön "
                      "életciklus. A tervdokumentum pedig megint más: az projektszintű referencia, "
                      "nem egy hibához tartozó fotó.")

    picture(prs, "ELSZÁMOLÁS", "TIG: a teljesítésigazolás előkészítése",
            "10-tig-teljesitesigazolas.png",
            caption="Elfogadott tételek csomagban, Excel és PDF export fotós melléklettel.",
            narration="A TIG a magyar építőiparban a számlázás alapja – enélkül nincs jogszerű számla. "
                      "Az app annyit csinál, hogy az elfogadott, fotóval bizonyított hibákat csomagba "
                      "zárja, és kiexportálja. Ez a modul köti össze a terepi munkát a pénzzel.")

    # ── 9-11. Leptetesek ────────────────────────────────────────────────────
    bullets(prs, "ÁLLAPOTGÉP", "Miért nem elég egy „kész” pipa?",
            ["Kilenc állapot, a nyitottól a lezártig – plusz a visszadobás.",
             "Nem lehet akárhonnan akárhová ugrani.",
             "A szerep is szűkít: az alvállalkozó nem fogadhatja el a saját munkáját.",
             "Ez adja az elszámolás bizonyíthatóságát."],
            narration="Itt a kulcsmondat: aki elvégzi a munkát, nem ugyanaz, mint aki igazolja. "
                      "Egy egyszerű 'kész' checkbox esetén az alvállalkozó saját magát minősítené. "
                      "Az állapotgép kényszeríti ki, hogy legyen egy ellenőrzési lépés a kész és az "
                      "elszámolható között.")

    picture(prs, "LÉPTETÉS", "A felület csak a megengedett lépéseket kínálja",
            "06-hiba-szerkesztes-allapotvaltas.png",
            caption="És amit a felület felkínál, azt a szerver még egyszer ellenőrzi.",
            narration="A legördülőben nem szerepel tiltott állapot. De ez önmagában nem védelem – "
                      "a böngészőben bármit át lehet írni. Ezért ugyanaz a szabály fut a szerveren is, "
                      "mentéskor. Ez visszatérő elv az egész appban: a felület kényelem, a szerver a fal.")

    bullets(prs, "VISSZALÉPTETÉS", "Visszafelé csak indokkal",
            ["Visszalépni csak három állapotból lehet.",
             "Mindegyikhez kötelező indokot írni – bekerül az eseménynaplóba.",
             "Az elfogadás visszavonása vezetői döntés."],
            lead="Egy elfogadott tétel visszavonása pénzügyi következménnyel jár, ezért nyomot hagy.",
            narration="Ez egy jó példa arra, hogyan lesz egy domain-szabályból kódszabály. Nem tiltjuk "
                      "meg a visszalépést – a valóságban szükség van rá, mert kiderülhet, hogy mégsem jó "
                      "a munka. De kötelezővé tesszük az indokot, hogy utólag ne lehessen csendben "
                      "átírni a történetet.")

    # ── 12-14. Jogosultsag ──────────────────────────────────────────────────
    bullets(prs, "JOGOSULTSÁG", "A szabály, amit be kellett tartatni",
            ["Adminisztrátor, projektvezető: minden projekt, minden hiba.",
             "Építésvezető, megtekintő: csak a saját projektjeik – abban viszont minden.",
             "Alvállalkozó: saját projekt ÉS saját cég. Egyik önmagában nem elég."],
            lead="Nem elég elrejteni a gombot – az adatot sem szabad kiadni.",
            narration="A kétdimenziós szűrés a nehéz rész. Két teszteset bizonyítja, hogy tényleg "
                      "működik: az egyik alvállalkozónak van saját cége hibája egy olyan projektben, "
                      "aminek nem tagja – azt nem láthatja. A másiknak van idegen cég hibája egy olyan "
                      "projektben, aminek tagja – azt sem láthatja. Ha csak az egyik tengelyen szűrnénk, "
                      "az egyik eset átcsúszna.")

    compare(prs, "UGYANAZ AZ OLDAL", "Amit a szerep eldönt",
            [("11-jogosultsag-admin.png", "Adminisztrátor – minden projekt"),
             ("12-jogosultsag-muvezeto.png", "Építésvezető – 1 projekt, minden cég"),
             ("13-jogosultsag-teszt1.png", "Alvállalkozó – 1 projekt, saját cég")],
            caption="Azonos URL, azonos kód. A különbséget kizárólag a bejelentkezett szerep adja.",
            narration="Ez a prezentáció legfontosabb képe – hagyj rajta időt. Ugyanaz az oldal, ugyanaz "
                      "a kód, három bejelentkezés. Balra kilenc hiba, jobbra egyetlen egy. Ezt "
                      "elmagyarázni sokáig tartana, megmutatni két másodperc.")

    bullets(prs, "NÉGY LÉPCSŐ", "Hogyan épült fel – és mi hiányzik még",
            ["1. Identitás: az Auth-fiók összekötése a profillal és a céggel.",
             "2. Írás: ki mit módosíthat – szerveroldalon kikényszerítve.",
             "3. Olvasás: ki mit lát – hatókör-szűrés minden listán.",
             "4. HÁTRAVAN: adatbázis-szintű védelem (RLS)."],
            lead="A Supabase publishable kulcs publikus – benne van a böngészőben.",
            narration="Ez az őszinte dia. Amíg nincs RLS, addig az alkalmazás szűr, de az adatbázis nem. "
                      "Aki megnyitja a fejlesztői eszközöket és közvetlenül hívja az API-t, megkerüli az "
                      "egészet. Ezért nem adtunk ki még valódi külsős fiókot. A tanulság általánosítható: "
                      "a jogosultság nem az, hogy 'úgyis csak az app hívja'.")

    # ── 15-17. Hogyan epult ─────────────────────────────────────────────────
    bullets(prs, "MÓDSZER", "Fejlesztési ritmus",
            ["Nem teljes specifikációból indultunk – rough inputokból.",
             "Séma → seed/mock adat → read-only UI → kontrollált írás.",
             "Minden kör végén lint, build és kézi teszt.",
             "Kis commitok, push csak jóváhagyás után."],
            narration="Ez a ritmus adta a biztonságot. A mock adat kulcsfontosságú volt: amíg a domain "
                      "modell változott, az app végigkattintható maradt. Aki elsőre az adatbázissal kezd, "
                      "az minden domain-döntésnél migrációt ír.")

    bullets(prs, "ITERÁCIÓ", "Két példa arra, hogyan nőtt a funkció",
            ["Fotós bizonyítás: metadata → Storage feltöltés → mobil viewer → törlés.",
             "Tervmérés: az AI-mennyiségszámítás helyett először egy kalibrálható mérő.",
             "A bizonytalan részt későbbre halasztottuk, a biztosat megépítettük."],
            lead="Egy „egyszerű” képfeltöltés valójában sok mobil UX részletből áll.",
            narration="A tervmérés jó példa a kockázatcsökkentésre. Az eredeti ötlet az volt, hogy az AI "
                      "olvassa le a tervrajzról a mennyiségeket. Ez nagy, bizonytalan feladat. Ehelyett "
                      "megépítettük a kalibrálható kézi mérőt: két pont megadja a méretarányt, onnantól "
                      "mérhető a terület és a hossz. Működik, használható, és ha később jön az AI, "
                      "erre épülhet.")

    compare(prs, "REAL-DEVICE FEEDBACK", "A valódi telefon döntött",
            [("15-mobil-attekintes.png", "Áttekintés"),
             ("16-mobil-hibalista.png", "Hibalista"),
             ("17-mobil-tobb-menu.png", "Navigáció")],
            caption="Sok hiba csak igazi eszközön jött elő – a dev szervert LAN-on futtatva teszteltük.",
            narration="A desktop böngésző mobil-emulátora nem elég. Ami csak telefonon jött elő: "
                      "hydration warningok, a fájlválasztó viselkedése, PDF-megjelenítési korlátok, "
                      "és a navigáció, amit többször újraterveztünk. A visszajelzés képernyőképekkel "
                      "érkezett, és sok kis körben javítottunk.")

    # ── 18-20. Uj tapasztalatok ─────────────────────────────────────────────
    bullets(prs, "AMIT A KÓD ELÁRULT", "Három rés, amit nem kézi teszt talált meg",
            ["A TIG export végpontjain SEM jogosultság-, SEM hatókör-ellenőrzés nem volt.",
             "A tervméréshez nem volt „fojtópont”: idegen projekt tervének adatai elérhetők voltak.",
             "Az akadály-módosítás megkerülte a hatókört."],
            lead="Módszeres kódátvizsgálás, nem kattintgatás.",
            narration="A TIG export volt a legsúlyosabb: bárki, aki be tudott lépni, letölthette bármelyik "
                      "csomag pénzügyi adatait – összegeket, alvállalkozói kapcsolattartót, fotókat. "
                      "Miért maradt rejtve? Mert mindig adminként próbáltuk, akinek amúgy is van joga hozzá. "
                      "Ez a kézi tesztelés vakfoltja: a saját jogosultságoddal nézed a rendszert.")

    bullets(prs, "AUTOMATIZÁLÁS", "43 end-to-end teszt, 55 másodperc",
            ["Hat teszt-fiók, valódi bejelentkezéssel – nem mockolt szerepekkel.",
             "A jogosultsági mátrix minden cellája külön állítás.",
             "A három javított rés külön őrizve.",
             "Visszatörtük a javítást: a teszt azonnal elbukott. Utána visszatettük."],
            lead="Egy zöld teszt nem bizonyíték. Az a bizonyíték, ha el TUD bukni.",
            narration="Az utolsó pont a legfontosabb. Írhatsz negyven tesztet, ami mindig zöld, és semmit "
                      "nem véd. Ezért szándékosan visszatörtük a TIG export javítását, és megnéztük, hogy "
                      "elbukik-e. Mind a négy korlátozott szerep tesztje azonnal piros lett. Csak ezután "
                      "hihetjük el, hogy a teszt valóban őrködik.")

    bullets(prs, "AMIT A TESZT TALÁLT", "Egy hiba, amit senki nem vett észre",
            ["A sikeres mentés visszajelzése a szerkesztő űrlapon BELÜL jelent meg.",
             "Az űrlap viszont mentéskor bezárul – az üzenet sosem látszott.",
             "Mindhárom részletező panelt érintette."],
            lead="A teszt nem csak azt fogja meg, amit keresel.",
            narration="Ezt hónapokig nem vette észre senki, pedig mindennap használtuk. Az ok emberi: "
                      "mentés után látod, hogy megváltozott az adat, és nem hiányolod a megerősítést. "
                      "A gépnek viszont feltűnt, mert azt az üzenetet kereste, aminek a kód szerint "
                      "meg kellett volna jelennie.")

    # ── 21-23. Zaras ────────────────────────────────────────────────────────
    bullets(prs, "COLLABORATION MODEL", "AI mint fejlesztőtárs, nem autopilot",
            ["Az AI bontott scope-ot, írt kódot, diagnosztizált és dokumentált.",
             "Az ember döntötte el: mi legyen most, mi maradjon későbbre.",
             "A jó prompt korlátokat, tesztet és commit-szabályt is tartalmazott.",
             "A legnagyobb érték nem a kódírás volt, hanem a módszeres átvizsgálás."],
            narration="Itt érdemes őszintének lenni: az AI sokszor magabiztosan tévedett is, és több "
                      "kört vett, mire egy-egy UI részlet jó lett. Ahol viszont egyértelműen többet "
                      "adott, mint amennyi időt elvett: a módszeres kódátvizsgálás. Három valódi "
                      "biztonsági rést talált, amit hónapok kézi tesztelése sem hozott elő.")

    bullets(prs, "HOL TARTUNK", "Állapot és a következő lépés",
            ["Használható MVP-váz több projektre, szerepekkel.",
             "Jogosultság 1–3. lépcső kész, automatikusan visszamérve.",
             "Következő: adatbázis-szintű RLS – enélkül nincs külsős fiók.",
             "A TIG üzleti folyamatot még pontosítani kell a szakértővel."],
            narration="A záró állapot. Ha egyetlen dolgot kell megjegyezni a következő lépésről: "
                      "az RLS nem opcionális finomítás, hanem az a pont, ami előtt nem szabad valódi "
                      "külsős felhasználót beengedni.")

    bullets(prs, "TAKEAWAYS", "Mit vigyél haza?",
            ["Kezdj vertical slice-szal, ne a teljes appal.",
             "A mock adat nem csalás, hanem gyors tanulási eszköz.",
             "Tesztelj azon az eszközön, amin használni fogják.",
             "Az elrejtett gomb nem biztonság – a végpontot kell védeni.",
             "Egy teszt akkor ér valamit, ha be tudod bizonyítani, hogy el tud bukni."],
            narration="Ha csak egy mondat marad meg: a felület kényelem, a szerver a fal. Minden más "
                      "ebből következik.")

    s = chrome(prs, "KÖSZÖNÖM", "Q&A")
    textbox(s, 1.33, 2.6, 15.0, 3.0,
            ["• Mi volt a legnagyobb tanulság?",
             "• Hol segített legtöbbet az AI, és hol nem?",
             "• Mit csinálnánk másképp a következő MVP-nél?"],
            25, BODY, spacing=18)
    note(s, "Lehetséges kérdések, amikre készülj: Mennyi idő volt? Mennyibe került? "
            "Mit csinálna másképp? Hogyan ellenőrzöd, hogy az AI nem hallucinál? "
            "(Erre jó válasz: lint, build, teszt, és a szándékos visszatörés-próba.)")

    # ── Oldalszamok ─────────────────────────────────────────────────────────
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        textbox(slide, 15.42, 9.31, 1.67, 0.2, f"{i} / {total}", 12, FAINT)

    try:
        prs.save(str(OUTPUT))
    except PermissionError:
        raise SystemExit(
            f"\nA fajl zarolva van, nem tudom felulirni:\n  {OUTPUT}\n\n"
            "Szinte biztosan meg van nyitva PowerPointban. Zard be, es futtasd ujra.\n"
        )

    with_notes = sum(1 for s in prs.slides if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip())
    print(f"Kesz: {OUTPUT}")
    print(f"  {total} dia, ebbol {with_notes} tartalmaz eloadoi jegyzetet")


if __name__ == "__main__":
    build()
